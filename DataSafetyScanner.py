#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
数据安全检查工具 Data Safety Scanner
=====================================
扫描本地文件中的个人敏感信息，支持查看详情、删除文件、脱敏处理。
完全离线运行，无需网络。

作者: 小宇
版本: 1.0
"""

import os
import re
import sys
import json
import shutil
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
from datetime import datetime
from pathlib import Path

# ============================================================
# 敏感信息规则配置
# ============================================================
SENSITIVE_PATTERNS = [
    {
        "name": "身份证号",
        "pattern": r'\b[1-9]\d{5}(?:19|20)\d{2}(?:0[1-9]|1[0-2])(?:0[1-9]|[12]\d|3[01])\d{3}[\dXx]\b',
        "description": "中国大陆18位身份证号码",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "手机号",
        "pattern": r'\b1[3-9]\d{9}\b',
        "description": "中国大陆手机号码",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "银行卡号",
        "pattern": r'\b(?:62|60|58|56|55|54|53|52|51|50|49|48|47|46|45|44|43|42|41|40|39|38|37|36|35|34|33|32|31|30)\d{14,17}\b',
        "description": "中国大陆银行卡号（16-19位）",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "电子邮箱",
        "pattern": r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b',
        "description": "电子邮箱地址",
        "risk": "中",
        "color": "#e67e22"
    },
    {
        "name": "固定电话",
        "pattern": r'\b(?:0\d{2,3}[-\s]?\d{7,8}|\(0\d{2,3}\)\s?\d{7,8})\b',
        "description": "中国大陆固定电话号码",
        "risk": "中",
        "color": "#e67e22"
    },
    {
        "name": "IP地址",
        "pattern": r'(?<!\d)(?:(?:25[0-5]|2[0-4]\d|[01]?\d\d?)\.){3}(?:25[0-5]|2[0-4]\d|[01]?\d\d?)(?!\d)',
        "description": "IPv4地址",
        "risk": "低",
        "color": "#f39c12"
    },
    {
        "name": "API密钥/Token",
        "pattern": r'(?:api[_-]?key|apikey|secret|token|password|passwd|pwd)\s*[:=]\s*["\']?[A-Za-z0-9_\-]{16,}["\']?',
        "description": "可能为API密钥、Token或密码（不区分大小写）",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "护照号",
        "pattern": r'\b[EeGg]\d{8}\b',
        "description": "中国大陆护照号码（E/G开头+8位数字）",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "详细地址（带标签）",
        "pattern": r'(?:地址|住址|家庭住址|户籍地址|通讯地址|联系地址|居住地|所在地)\s*[：:]\s*\S{4,}',
        "description": "带有地址标签的完整地址，如\"地址：广州市天河区...\"",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "四级地址（省市区路号）",
        "pattern": r'(?:[^\d_]{2,}(?:省|自治区))[^\d_]{1,}(?:市)[^\d_]{1,}(?:区|县|镇)[^\d_]{1,}(?:路|街|巷|道|大道|大街)\S*?(?:号|巷|弄)',
        "description": "省+市+区+路+号的完整地址",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "三级地址（省市组合）",
        "pattern": r'(?:(?:[^\d_]{2,}(?:省|自治区)[^\d_]{1,}(?:市|区|县))|(?:[^\d_]{2,}(?:市)[^\d_]{1,}(?:区|县)[^\d_]{1,}(?:路|街|巷|道)))',
        "description": "省+市或市+区+路的组合地址",
        "risk": "中",
        "color": "#e67e22"
    },
    {
        "name": "具体门牌号码",
        "pattern": r'(?:(?:路|街|巷|道|大道|大街|弄)\s*[0-9]{1,5}\s*(?:号|弄|巷|室))|(?:[0-9]{1,5}\s*栋\s*(?:[0-9]{1,2}\s*单元\s*[0-9]{1,4}\s*室|[0-9]{1,4}\s*室))',
        "description": "具体门牌号+栋+单元+室组合，如\"科技路88号\"或\"12栋3单元401室\"",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "证件号/账户",
        "pattern": r'(?:证件号|证件号码|账户(?:号|号码)|账号[：:]|帐号[：:])\s*\S{4,}',
        "description": "带标签的证件号码或账户标识，如\"证件号：110101...\"",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "密码/验证码",
        "pattern": r'(?:密码|验证码|校验码|动态码|短信验证码)\s*[：:]\s*\S{4,}',
        "description": "带标签的密码或验证码，如\"密码：Abc123...\"",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "定位/轨迹",
        "pattern": r'(?:定位[：:]|地理位置[：:]|GPS坐标[：:]|经纬度[：:]|行动轨迹[：:])\s*\S{2,}',
        "description": "带标签的定位或轨迹信息，如\"定位：113.32,23.12\"",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "生物识别信息",
        "pattern": r'(?:人脸信息|指纹信息|虹膜信息|声纹信息|掌纹信息|面部识别信息|生物特征信息|生物识别信息)',
        "description": "明确标注的生物识别信息，如\"指纹信息\"",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "医疗/病历",
        "pattern": r'(?:病历[：：]|病史[：：]|就诊记录[：：]|诊断结果[：：]|检查报告[：：])|(?:病历信息|病史信息|就诊记录|住院记录|检查报告|医疗记录|健康档案)(?:\s*[：:]\s*\S{2,}|$)',
        "description": "病历、就诊记录等健康信息，需带标签或明确标记",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "征信/贷款/流水",
        "pattern": r'(?:征信报告|信用报告|交易流水|银行流水|流水记录|收支明细)(?:\s*(?:[：:]\s*\S{2,}|信息|数据|记录|单)|$)',
        "description": "征信报告、银行流水等金融文件，排除普通\"贷款\"一词",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "薪资/收入",
        "pattern": r'(?:工资条|工资单|薪资[：:]\s*\S{2,}|薪酬[：:]\s*\S{2,}|待遇[：:]\s*\S{2,})',
        "description": "工资条/工资单，或带标签的薪资信息，排除普通\"收入\"一词",
        "risk": "高",
        "color": "#e74c3c"
    },
    {
        "name": "未成年人信息",
        "pattern": r'(?:未成年人|未成年)\s*(?:信息|数据|资料|名单|隐私)|(?:儿童|学生|幼儿|青少年)\s*(?:个人信息|隐私信息|隐私数据|名单)',
        "description": "未成年人个人信息，如\"未成年人信息\"或\"儿童隐私数据\"",
        "risk": "高",
        "color": "#e74c3c"
    },
]

# 支持的扩展名及读取方式
SUPPORTED_EXTENSIONS = {
    '.txt': 'text',
    '.md': 'text',
    '.csv': 'text',
    '.log': 'text',
    '.json': 'text',
    '.xml': 'text',
    '.yaml': 'text',
    '.yml': 'text',
    '.ini': 'text',
    '.cfg': 'text',
    '.conf': 'text',
    '.sql': 'text',
    '.html': 'text',
    '.htm': 'text',
    '.css': 'text',
    '.js': 'text',
    '.py': 'text',
    '.java': 'text',
    '.c': 'text',
    '.cpp': 'text',
    '.h': 'text',
    '.sh': 'text',
    '.bat': 'text',
    '.ps1': 'text',
    '.env': 'text',
    '.properties': 'text',
    '.xlsx': 'excel',
    '.xls': 'excel',
    '.docx': 'word',
    '.doc': 'word',
    '.pdf': 'pdf',
    '.pptx': 'ppt',
    '.ppt': 'ppt',
    '.sqlite': 'database',
    '.db': 'database',
    '.db3': 'database',
    '.sqlite3': 'database',
}

# 最大文件大小（默认10MB）
MAX_FILE_SIZE = 10 * 1024 * 1024

# ============================================================
# 核心扫描引擎
# ============================================================
class ScanEngine:
    """扫描引擎 - 在单独线程中运行"""

    def __init__(self):
        self.stop_flag = threading.Event()
        self.scanned_files = 0
        self.matched_files = 0
        self.total_matches = 0

    def stop(self):
        self.stop_flag.set()

    def scan_file(self, file_path, enabled_patterns):
        """扫描单个文件，返回匹配结果列表"""
        if self.stop_flag.is_set():
            return []

        ext = os.path.splitext(file_path)[1].lower()
        if ext not in SUPPORTED_EXTENSIONS:
            return []

        # 检查文件大小
        try:
            file_size = os.path.getsize(file_path)
            if file_size > MAX_FILE_SIZE:
                return []  # 跳过过大文件
            if file_size == 0:
                return []
        except (OSError, IOError):
            return []

        # Excel文件处理
        if SUPPORTED_EXTENSIONS.get(ext) == 'excel':
            return self._scan_excel_file(file_path, enabled_patterns)

        # Word文档处理
        if SUPPORTED_EXTENSIONS.get(ext) == 'word':
            return self._scan_word_file(file_path, enabled_patterns)

        # PDF文档处理
        if SUPPORTED_EXTENSIONS.get(ext) == 'pdf':
            return self._scan_pdf_file(file_path, enabled_patterns)

        # PPT演示文稿处理
        if SUPPORTED_EXTENSIONS.get(ext) == 'ppt':
            return self._scan_ppt_file(file_path, enabled_patterns)

        # 数据库文件处理
        if SUPPORTED_EXTENSIONS.get(ext) == 'database':
            return self._scan_database_file(file_path, enabled_patterns)

        # 尝试以文本方式读取（纯文本文件）
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
        except Exception:
            try:
                with open(file_path, 'r', encoding='gbk', errors='replace') as f:
                    content = f.read()
            except Exception:
                return []

        matches = []
        lines = content.split('\n')

        for pattern_info in enabled_patterns:
            try:
                regex = re.compile(pattern_info['pattern'])
            except re.error:
                continue

            for line_idx, line in enumerate(lines):
                if self.stop_flag.is_set():
                    return matches

                for match in regex.finditer(line):
                    matched_text = match.group()
                    start_pos = match.start()

                    # 获取上下文（前后各20个字符）
                    context_start = max(0, start_pos - 20)
                    context_end = min(len(line), start_pos + len(matched_text) + 20)
                    context = line[context_start:context_end].strip()

                    # 关键词地址类型特殊处理 - 需要更多上下文判断
                    if pattern_info['name'] == '住址关键词':
                        # 取整行作为上下文
                        context = line.strip()
                        if len(context) < 10:  # 太短可能只是无关匹配
                            continue

                    matches.append({
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'file_dir': os.path.dirname(file_path),
                        'line_number': line_idx + 1,
                        'matched_text': matched_text[:100],  # 截断过长匹配
                        'pattern_name': pattern_info['name'],
                        'risk': pattern_info['risk'],
                        'context': context,
                        'full_line': line.strip(),
                    })

        return matches

    def _scan_excel_file(self, file_path, enabled_patterns):
        """扫描Excel文件（.xlsx / .xls），读取每个单元格的内容"""
        if self.stop_flag.is_set():
            return []
        matches = []
        try:
            ext = os.path.splitext(file_path)[1].lower()
            if ext == '.xlsx':
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            else:
                import xlrd
                wb = xlrd.open_workbook(file_path)
        except Exception:
            return []

        for sheet_name in wb.sheet_names:
            if self.stop_flag.is_set():
                break
            try:
                if ext == '.xlsx':
                    ws = wb[sheet_name]
                    for row_idx, row in enumerate(ws.iter_rows(values_only=True), 1):
                        self._scan_excel_row(row, row_idx, sheet_name, file_path, enabled_patterns, matches)
                else:
                    ws = wb.sheet_by_name(sheet_name)
                    for row_idx in range(ws.nrows):
                        row = ws.row_values(row_idx)
                        self._scan_excel_row(row, row_idx + 1, sheet_name, file_path, enabled_patterns, matches)
            except Exception:
                continue

        if ext == '.xlsx':
            wb.close()
        return matches

    def _scan_excel_row(self, row, row_idx, sheet_name, file_path, enabled_patterns, matches):
        """扫描Excel的一行数据"""
        for col_idx, cell_value in enumerate(row):
            if self.stop_flag.is_set():
                return
            if cell_value is None:
                continue
            cell_text = str(cell_value).strip()
            if not cell_text or len(cell_text) < 2:
                continue

            for pattern_info in enabled_patterns:
                try:
                    regex = re.compile(pattern_info['pattern'])
                except re.error:
                    continue
                for match in regex.finditer(cell_text):
                    matched_text = match.group()
                    start_pos = match.start()
                    context_start = max(0, start_pos - 20)
                    context_end = min(len(cell_text), start_pos + len(matched_text) + 20)
                    context = cell_text[context_start:context_end].strip()

                    # 列号转字母（A, B, ..., Z, AA, AB...）
                    col_letter = ''
                    c = col_idx
                    while c >= 0:
                        col_letter = chr(65 + (c % 26)) + col_letter
                        c = c // 26 - 1
                        if c < 0:
                            break
                    cell_ref = f"{col_letter}{row_idx}"

                    matches.append({
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'file_dir': os.path.dirname(file_path),
                        'line_number': row_idx,
                        'matched_text': matched_text[:100],
                        'pattern_name': pattern_info['name'],
                        'risk': pattern_info['risk'],
                        'context': context,
                        'full_line': cell_text,
                        'cell_ref': cell_ref,
                        'sheet_name': sheet_name,
                    })

    def _scan_word_file(self, file_path, enabled_patterns):
        """扫描Word文档（.docx），读取段落和表格中的文本"""
        if self.stop_flag.is_set():
            return []
        try:
            from docx import Document
            doc = Document(file_path)
        except Exception:
            return []

        matches = []
        # 扫描段落
        for para_idx, para in enumerate(doc.paragraphs):
            if self.stop_flag.is_set():
                break
            text = para.text.strip()
            if not text or len(text) < 2:
                continue
            matches.extend(self._scan_text_block(text, para_idx + 1, file_path, enabled_patterns, 'paragraph'))

        # 扫描表格
        for table_idx, table in enumerate(doc.tables):
            if self.stop_flag.is_set():
                break
            for row_idx, row in enumerate(table.rows):
                for cell in row.cells:
                    text = cell.text.strip()
                    if not text or len(text) < 2:
                        continue
                    line_num = f"table{table_idx + 1}_row{row_idx + 1}"
                    matches.extend(self._scan_text_block(text, line_num, file_path, enabled_patterns, 'table'))
        return matches

    def _scan_pdf_file(self, file_path, enabled_patterns):
        """扫描PDF文件，提取每页文本"""
        if self.stop_flag.is_set():
            return []
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
        except Exception:
            return []

        matches = []
        for page_idx, page in enumerate(reader.pages):
            if self.stop_flag.is_set():
                break
            text = page.extract_text()
            if not text:
                continue
            lines = text.split('\n')
            for line_idx, line in enumerate(lines):
                if self.stop_flag.is_set():
                    break
                line = line.strip()
                if not line or len(line) < 2:
                    continue
                line_num = f"p{page_idx + 1}_l{line_idx + 1}"
                matches.extend(self._scan_text_block(line, line_num, file_path, enabled_patterns, 'page'))
        return matches

    def _scan_ppt_file(self, file_path, enabled_patterns):
        """扫描PowerPoint演示文稿（.pptx），读取所有幻灯片中的文本"""
        if self.stop_flag.is_set():
            return []
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
        except Exception:
            return []

        matches = []
        for slide_idx, slide in enumerate(prs.slides):
            if self.stop_flag.is_set():
                break
            for shape in slide.shapes:
                if self.stop_flag.is_set():
                    break
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text or len(text) < 2:
                            continue
                        line_num = f"s{slide_idx + 1}"
                        matches.extend(self._scan_text_block(text, line_num, file_path, enabled_patterns, 'slide'))
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if not text or len(text) < 2:
                                continue
                            line_num = f"s{slide_idx + 1}"
                            matches.extend(self._scan_text_block(text, line_num, file_path, enabled_patterns, 'slide'))
        return matches

    def _scan_database_file(self, file_path, enabled_patterns):
        """扫描SQLite数据库文件，读取所有表的文本字段"""
        if self.stop_flag.is_set():
            return []
        try:
            import sqlite3
            conn = sqlite3.connect(file_path)
            conn.text_factory = str
            cursor = conn.cursor()
        except Exception:
            return []

        matches = []
        try:
            # 获取所有表名
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
            tables = [row[0] for row in cursor.fetchall()]

            for table_name in tables:
                if self.stop_flag.is_set():
                    break
                try:
                    cursor.execute(f"PRAGMA table_info(\"{table_name}\")")
                    columns = [row[1] for row in cursor.fetchall()]

                    # 只处理文本类型列
                    text_cols = []
                    for col in columns:
                        cursor.execute(f"PRAGMA table_info(\"{table_name}\")")
                        col_info = [r for r in cursor.fetchall() if r[1] == col]
                        if col_info and col_info[0][2].upper() in ('TEXT', 'VARCHAR', 'CHAR', 'CLOB', 'NVARCHAR', 'NCHAR', ''):
                            text_cols.append(col)

                    if not text_cols:
                        continue

                    # 读取数据（最多前10000行）
                    cursor.execute(f"SELECT rowid, * FROM \"{table_name}\" LIMIT 10000")
                    for row in cursor.fetchall():
                        if self.stop_flag.is_set():
                            break
                        rowid = row[0]
                        for col_idx, col_name in enumerate(text_cols, 1):
                            cell_value = row[col_idx] if col_idx < len(row) else None
                            if cell_value is None:
                                continue
                            cell_text = str(cell_value).strip()
                            if not cell_text or len(cell_text) < 2:
                                continue
                            matches.extend(self._scan_text_block(
                                cell_text,
                                f"row{rowid}",
                                file_path, enabled_patterns, 'database'
                            ))
                            # 批量附加表名和列名信息
                            for m in matches:
                                if m['source_type'] == 'database':
                                    m['db_table'] = table_name
                                    m['db_column'] = col_name
                except Exception:
                    continue
        finally:
            conn.close()
        return matches

    def _scan_text_block(self, text, line_number, file_path, enabled_patterns, source_type):
        """通用文本块扫描逻辑，供Word/PDF共用"""
        if self.stop_flag.is_set():
            return []
        results = []
        for pattern_info in enabled_patterns:
            if self.stop_flag.is_set():
                break
            try:
                regex = re.compile(pattern_info['pattern'])
            except re.error:
                continue
            for match in regex.finditer(text):
                matched_text = match.group()
                start_pos = match.start()
                context_start = max(0, start_pos - 20)
                context_end = min(len(text), start_pos + len(matched_text) + 20)
                context = text[context_start:context_end].strip()
                results.append({
                    'file_path': file_path,
                    'file_name': os.path.basename(file_path),
                    'file_dir': os.path.dirname(file_path),
                    'line_number': line_number,
                    'matched_text': matched_text[:100],
                    'pattern_name': pattern_info['name'],
                    'risk': pattern_info['risk'],
                    'context': context,
                    'full_line': text,
                    'source_type': source_type,
                })
        return results

    def scan_directory(self, directory, enabled_patterns, progress_callback=None):
        """递归扫描目录"""
        results = []
        self.scanned_files = 0
        self.matched_files = 0
        self.total_matches = 0

        for root, dirs, files in os.walk(directory):
            # 跳过隐藏目录和系统目录
            dirs[:] = [d for d in dirs if not d.startswith('.') and d not in
                       ('node_modules', '__pycache__', '.git', '.svn', 'venv', '.venv', 'env')]

            if self.stop_flag.is_set():
                break

            for file in files:
                if self.stop_flag.is_set():
                    break

                file_path = os.path.join(root, file)
                self.scanned_files += 1

                file_matches = self.scan_file(file_path, enabled_patterns)
                if file_matches:
                    self.matched_files += 1
                    self.total_matches += len(file_matches)
                    results.extend(file_matches)

                if progress_callback:
                    progress_callback(self.scanned_files, self.matched_files, file_path)

        return results


# ============================================================
# 操作工具函数
# ============================================================
def redact_text(text, pattern):
    """将匹配的敏感信息替换为***"""
    try:
        regex = re.compile(pattern)
        return regex.sub(lambda m: '*' * len(m.group()), text)
    except re.error:
        return text


def redact_file(file_path, pattern, pattern_name):
    """对文件中匹配的内容进行脱敏处理（创建备份后替换）"""
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return False, "不支持的文件格式"

    # 创建备份
    bak_path = file_path + '.bak'
    try:
        shutil.copy2(file_path, bak_path)
    except Exception as e:
        return False, f"备份失败: {str(e)}"

    # 读取并脱敏
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
    except Exception:
        try:
            with open(file_path, 'r', encoding='gbk', errors='replace') as f:
                content = f.read()
        except Exception as e:
            os.remove(bak_path)
            return False, f"读取失败: {str(e)}"

    original_content = content
    content = redact_text(content, pattern)

    if content == original_content:
        os.remove(bak_path)
        return False, "未找到匹配内容"

    try:
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(content)
        return True, f"已脱敏，备份文件: {bak_path}"
    except Exception as e:
        # 恢复备份
        shutil.copy2(bak_path, file_path)
        os.remove(bak_path)
        return False, f"写入失败: {str(e)}"


def safe_delete_file(file_path):
    """安全删除文件（移到回收站/Trash）"""
    import subprocess
    try:
        # macOS先移到Trash
        subprocess.run(['osascript', '-e',
            f'tell app "Finder" to delete (POSIX file "{file_path}" as alias)'],
            capture_output=True, timeout=10)
        return True, "文件已移至废纸篓"
    except Exception:
        # 后备方案：重命名后删除
        try:
            trash_name = file_path + '.deleted'
            os.rename(file_path, trash_name)
            os.remove(trash_name)
            return True, "文件已删除"
        except Exception as e:
            return False, f"删除失败: {str(e)}"


# ============================================================
# Tkinter GUI应用
# ============================================================
class DataSafetyScannerApp:
    def __init__(self, root):
        self.root = root
        self.root.title("◆ 钻石专属 ◆ 数据安全检查工具 v1.0 - 开平农商银行·数据银行部")
        self.root.geometry("1200x750")
        self.root.minsize(900, 600)

        # 设置主题色
        self.bg_color = "#f5f6fa"
        self.accent_color = "#2c3e50"
        self.danger_color = "#e74c3c"
        self.warning_color = "#e67e22"
        self.success_color = "#27ae60"

        try:
            self.root.iconbitmap(default=None)  # macOS下忽略图标
        except Exception:
            pass

        # 状态变量
        self.scan_results = []  # 所有扫描结果
        self.filtered_results = []  # 当前过滤后的结果
        self.is_scanning = False
        self.scan_thread = None
        self.engine = ScanEngine()
        self.selected_patterns = [p for p in SENSITIVE_PATTERNS]  # 默认全选

        # 扫描目录列表
        self.scan_dirs = []

        self._build_ui()
        self._setup_styles()

    def _setup_styles(self):
        style = ttk.Style()
        style.theme_use('clam')
        style.configure('Treeview', rowheight=28, font=('Helvetica', 12))
        style.configure('Treeview.Heading', font=('Helvetica', 13, 'bold'), padding=(5, 5))
        style.configure('High.Treeview', foreground='#e74c3c')
        style.configure('Medium.Treeview', foreground='#e67e22')
        style.configure('Low.Treeview', foreground='#f39c12')
        style.configure('TButton', font=('Helvetica', 12), padding=(10, 5))
        style.configure('Header.TLabel', font=('Helvetica', 16, 'bold'), foreground=self.accent_color)
        style.configure('Status.TLabel', font=('Helvetica', 11))

    def _build_ui(self):
        # 主容器
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)

        # ===== 顶部标题栏 =====
        header_frame = ttk.Frame(main_frame)
        header_frame.pack(fill=tk.X, pady=(0, 5))

        title_label = ttk.Label(header_frame, text="🔍 数据安全检查工具", style='Header.TLabel')
        title_label.pack(side=tk.LEFT)

        self.status_label = ttk.Label(header_frame, text="就绪", style='Status.TLabel')
        self.status_label.pack(side=tk.RIGHT, padx=10)

        # 钻石专属标志横幅
        badge_frame = ttk.Frame(main_frame)
        badge_frame.pack(fill=tk.X, pady=(0, 10))
        badge_label = tk.Label(badge_frame,
                               text=" ◆ 钻石专属 ◆  广东省·开平农商银行·数据银行部  ◆ 钻石专属 ◆ ",
                               font=('Helvetica', 10), fg="#b8860b", bg="#fff8dc",
                               relief=tk.GROOVE, bd=1, padx=5, pady=3)
        badge_label.pack(fill=tk.X)

        # ===== 配置区域 =====
        config_frame = ttk.LabelFrame(main_frame, text="扫描配置", padding="10")
        config_frame.pack(fill=tk.X, pady=(0, 10))

        # 第一行：目录选择
        dir_frame = ttk.Frame(config_frame)
        dir_frame.pack(fill=tk.X, pady=(0, 8))

        ttk.Label(dir_frame, text="扫描目录:", font=('Helvetica', 12)).pack(side=tk.LEFT)

        self.dir_listbox = tk.Listbox(dir_frame, height=3, font=('Helvetica', 11))
        self.dir_listbox.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(10, 5))

        dir_btn_frame = ttk.Frame(dir_frame)
        dir_btn_frame.pack(side=tk.RIGHT)

        ttk.Button(dir_btn_frame, text="添加目录", command=self._add_directory).pack(side=tk.TOP, pady=1)
        ttk.Button(dir_btn_frame, text="移除选中", command=self._remove_directory).pack(side=tk.TOP, pady=1)

        # 第二行：敏感信息类型选择
        pattern_frame = ttk.LabelFrame(config_frame, text="敏感信息类型", padding="5")
        pattern_frame.pack(fill=tk.X, pady=(0, 5))

        self.pattern_vars = {}
        col_count = 0
        for i, p in enumerate(SENSITIVE_PATTERNS):
            var = tk.BooleanVar(value=True)
            self.pattern_vars[p['name']] = var
            risk_color = p['color']
            cb = ttk.Checkbutton(pattern_frame, text=f"{p['name']} ({p['risk']})",
                                 variable=var)
            cb.grid(row=i // 4, column=i % 4, sticky='w', padx=10, pady=2)

        # 更多选项
        options_frame = ttk.Frame(config_frame)
        options_frame.pack(fill=tk.X)

        self.skip_large_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(options_frame, text="跳过大于10MB的文件", variable=self.skip_large_var).pack(side=tk.LEFT, padx=10)

        # 操作按钮
        action_frame = ttk.Frame(config_frame)
        action_frame.pack(fill=tk.X, pady=(8, 0))

        self.start_btn = ttk.Button(action_frame, text="▶ 开始扫描", command=self._start_scan)
        self.start_btn.pack(side=tk.LEFT, padx=(0, 10))

        self.stop_btn = ttk.Button(action_frame, text="⏹ 停止扫描", command=self._stop_scan, state=tk.DISABLED)
        self.stop_btn.pack(side=tk.LEFT, padx=(0, 10))

        ttk.Button(action_frame, text="🗑 清空结果", command=self._clear_results).pack(side=tk.LEFT, padx=(0, 10))

        # 导出报告
        ttk.Button(action_frame, text="📄 导出报告", command=self._export_report).pack(side=tk.RIGHT)

        # 过滤
        filter_frame = ttk.Frame(action_frame)
        filter_frame.pack(side=tk.RIGHT, padx=20)

        ttk.Label(filter_frame, text="风险过滤:").pack(side=tk.LEFT)
        self.filter_var = tk.StringVar(value="全部")
        filter_combo = ttk.Combobox(filter_frame, textvariable=self.filter_var,
                                    values=["全部", "高风险", "中风险", "低风险"],
                                    state="readonly", width=10)
        filter_combo.pack(side=tk.LEFT, padx=5)
        filter_combo.bind('<<ComboboxSelected>>', lambda e: self._apply_filter())

        # ===== 扫描进度 =====
        progress_frame = ttk.Frame(main_frame)
        progress_frame.pack(fill=tk.X, pady=(0, 5))

        self.progress_bar = ttk.Progressbar(progress_frame, mode='determinate')
        self.progress_bar.pack(side=tk.LEFT, fill=tk.X, expand=True)

        self.progress_label = ttk.Label(progress_frame, text="", width=30)
        self.progress_label.pack(side=tk.RIGHT, padx=10)

        # ===== 结果表格 =====
        result_frame = ttk.LabelFrame(main_frame, text="扫描结果", padding="5")
        result_frame.pack(fill=tk.BOTH, expand=True)

        # 创建表格
        columns = ('file', 'type', 'risk', 'match', 'line', 'dir')
        self.tree = ttk.Treeview(result_frame, columns=columns, show='headings',
                                  selectmode='extended')

        self.tree.heading('file', text='文件名')
        self.tree.heading('type', text='敏感信息类型')
        self.tree.heading('risk', text='风险')
        self.tree.heading('match', text='匹配内容')
        self.tree.heading('line', text='行号')
        self.tree.heading('dir', text='所在目录')

        self.tree.column('file', width=180, minwidth=100)
        self.tree.column('type', width=120, minwidth=80)
        self.tree.column('risk', width=60, minwidth=50)
        self.tree.column('match', width=300, minwidth=150)
        self.tree.column('line', width=60, minwidth=50)
        self.tree.column('dir', width=400, minwidth=200)

        # 滚动条
        scroll_y = ttk.Scrollbar(result_frame, orient=tk.VERTICAL, command=self.tree.yview)
        scroll_x = ttk.Scrollbar(result_frame, orient=tk.HORIZONTAL, command=self.tree.xview)
        self.tree.configure(yscrollcommand=scroll_y.set, xscrollcommand=scroll_x.set)

        self.tree.grid(row=0, column=0, sticky='nsew')
        scroll_y.grid(row=0, column=1, sticky='ns')
        scroll_x.grid(row=1, column=0, sticky='ew')

        result_frame.grid_rowconfigure(0, weight=1)
        result_frame.grid_columnconfigure(0, weight=1)

        # 绑定右键菜单和双击事件
        self.tree.bind('<Button-3>', self._show_context_menu)
        self.tree.bind('<Control-Button-1>', self._show_context_menu)  # macOS右键
        self.tree.bind('<Double-1>', self._on_double_click)

        # ===== 底部统计 =====
        stats_frame = ttk.Frame(main_frame)
        stats_frame.pack(fill=tk.X, pady=(5, 0))

        self.stats_label = ttk.Label(stats_frame, text="已扫描: 0 个文件 | 发现: 0 个文件有风险 | 共计: 0 条匹配",
                                      font=('Helvetica', 11))
        self.stats_label.pack(side=tk.LEFT)

        self.selected_info_label = ttk.Label(stats_frame, text="", font=('Helvetica', 11))
        self.selected_info_label.pack(side=tk.RIGHT)

        # 底部版权栏
        copyright_frame = ttk.Frame(main_frame)
        copyright_frame.pack(fill=tk.X)
        ttk.Label(copyright_frame, text="◆ 钻石专属 ◆  © 开平农商银行·数据银行部",
                  font=('Helvetica', 9), foreground="#999").pack(side=tk.RIGHT, padx=5)

        # ===== 右键菜单 =====
        self.context_menu = tk.Menu(self.root, tearoff=0, font=('Helvetica', 12))
        self.context_menu.add_command(label="📂 打开文件所在位置", command=self._open_file_location)
        self.context_menu.add_command(label="📄 打开文件", command=self._open_selected_file)
        self.context_menu.add_separator()
        self.context_menu.add_command(label="✂️ 脱敏处理（替换为***）", command=self._redact_selected)
        self.context_menu.add_command(label="🗑 删除文件（移入废纸篓）", command=self._delete_selected)
        self.context_menu.add_separator()
        self.context_menu.add_command(label="📋 复制匹配内容", command=self._copy_match_content)

    # ==================== 事件处理 ====================

    def _add_directory(self):
        dir_path = filedialog.askdirectory(title="选择要扫描的目录")
        if dir_path and dir_path not in self.scan_dirs:
            self.scan_dirs.append(dir_path)
            self.dir_listbox.insert(tk.END, dir_path)

    def _remove_directory(self):
        selection = self.dir_listbox.curselection()
        if selection:
            index = selection[0]
            self.scan_dirs.pop(index)
            self.dir_listbox.delete(index)

    def _get_enabled_patterns(self):
        enabled = []
        for p in SENSITIVE_PATTERNS:
            if self.pattern_vars[p['name']].get():
                enabled.append(p)
        return enabled

    def _start_scan(self):
        if not self.scan_dirs:
            messagebox.showwarning("提示", "请先添加要扫描的目录")
            return

        enabled = self._get_enabled_patterns()
        if not enabled:
            messagebox.showwarning("提示", "请至少选择一种敏感信息类型")
            return

        # 清空之前结果
        self._clear_results()
        self.is_scanning = True
        self.start_btn.config(state=tk.DISABLED)
        self.stop_btn.config(state=tk.NORMAL)
        self.status_label.config(text="扫描中...", foreground=self.warning_color)
        self.progress_bar['value'] = 0

        # 在单独线程中执行扫描
        dirs_to_scan = list(self.scan_dirs)
        self.engine = ScanEngine()

        def scan_thread_func():
            all_results = []
            total_dirs = len(dirs_to_scan)

            for dir_idx, directory in enumerate(dirs_to_scan):
                if self.engine.stop_flag.is_set():
                    break

                def make_progress(dir_idx, total_dirs):
                    def progress(scanned, matched, current_file):
                        self.root.after(0, lambda s=scanned, m=matched, f=current_file:
                            self._update_progress(s, m, f, dir_idx, total_dirs))
                    return progress

                results = self.engine.scan_directory(
                    directory, enabled,
                    progress_callback=make_progress(dir_idx, total_dirs)
                )
                all_results.extend(results)

            self.root.after(0, lambda r=all_results: self._scan_complete(r))

        self.scan_thread = threading.Thread(target=scan_thread_func, daemon=True)
        self.scan_thread.start()

    def _update_progress(self, scanned, matched, current_file):
        self.progress_label.config(text=f"扫描: {scanned} 个文件")
        # 进度条动画（不确定总量，用脉冲模式）
        self.progress_bar['mode'] = 'indeterminate'
        self.progress_bar.step(1)

        # 显示当前正在扫描的文件（简短显示）
        short_path = current_file if len(current_file) < 80 else '...' + current_file[-77:]
        self.status_label.config(text=f"正在扫描: {short_path}")

    def _scan_complete(self, results):
        self.is_scanning = False
        self.start_btn.config(state=tk.NORMAL)
        self.stop_btn.config(state=tk.DISABLED)
        self.progress_bar['mode'] = 'determinate'
        self.progress_bar['value'] = 100

        self.scan_results = results
        self._apply_filter()

        # 更新统计
        unique_files = len(set(r['file_path'] for r in results))
        self.stats_label.config(
            text=f"已扫描: {self.engine.scanned_files} 个文件 | 发现: {unique_files} 个文件有风险 | 共计: {len(results)} 条匹配"
        )

        if results:
            self.status_label.config(
                text=f"扫描完成！发现 {len(results)} 条敏感信息匹配",
                foreground=self.danger_color
            )
            messagebox.showinfo("扫描完成",
                f"扫描完成！\n\n"
                f"📁 扫描文件: {self.engine.scanned_files} 个\n"
                f"⚠️ 发现风险文件: {unique_files} 个\n"
                f"🔍 匹配条数: {len(results)} 条\n\n"
                f"建议逐一检查并处理高风险匹配项。")
        else:
            self.status_label.config(text="扫描完成，未发现敏感信息 ✓", foreground=self.success_color)
            messagebox.showinfo("扫描完成", "扫描完成，未发现敏感信息 ✓")

    def _stop_scan(self):
        if self.engine:
            self.engine.stop()
        self.status_label.config(text="正在停止...")
        self.stop_btn.config(state=tk.DISABLED)

    def _clear_results(self):
        self.scan_results = []
        self.filtered_results = []
        for item in self.tree.get_children():
            self.tree.delete(item)
        self.stats_label.config(text="已扫描: 0 个文件 | 发现: 0 个文件有风险 | 共计: 0 条匹配")
        self.status_label.config(text="就绪", foreground='black')

    def _apply_filter(self):
        filter_value = self.filter_var.get()
        risk_map = {"高风险": "高", "中风险": "中", "低风险": "低"}

        if filter_value == "全部":
            self.filtered_results = list(self.scan_results)
        else:
            risk_level = risk_map.get(filter_value, "")
            self.filtered_results = [r for r in self.scan_results if r['risk'] == risk_level]

        # 刷新表格
        for item in self.tree.get_children():
            self.tree.delete(item)

        for r in self.filtered_results:
            file_rel = r['file_name']
            # 根据文件类型显示位置信息
            if r.get('cell_ref'):
                line_display = f"{r['sheet_name']}!{r['cell_ref']}"
            elif r.get('source_type') == 'page':
                # PDF: p1_l5 → 第1页 第5行
                parts = str(r['line_number']).split('_l')
                page = parts[0][1:]  # 'p1' → '1'
                line = parts[1] if len(parts) > 1 else '?'
                line_display = f"第{page}页 第{line}行"
            elif r.get('source_type') == 'slide':
                line_display = f"第{r['line_number'][1:]}张幻灯片"
            elif r.get('source_type') == 'database':
                line_display = f"{r.get('db_table','?')}.{r.get('db_column','?')} #{r['line_number'][3:]}"
            elif isinstance(r['line_number'], int):
                line_display = f"第 {r['line_number']} 行"
            else:
                line_display = str(r['line_number'])
            self.tree.insert('', tk.END, values=(
                file_rel,
                r['pattern_name'],
                r['risk'],
                r['matched_text'][:80],
                line_display,
                r['file_dir']
            ))

        # 更新选中信息
        total = len(self.filtered_results)
        if total != len(self.scan_results):
            self.selected_info_label.config(text=f"（当前显示: {total}/{len(self.scan_results)} 条）")
        else:
            self.selected_info_label.config(text="")

    def _show_context_menu(self, event):
        selection = self.tree.selection()
        if selection:
            self.context_menu.post(event.x_root, event.y_root)

    def _get_selected_results(self):
        """获取当前选中行对应的结果"""
        selection = self.tree.selection()
        selected = []
        for item_id in selection:
            values = self.tree.item(item_id, 'values')
            # 通过文件名+行号+匹配内容定位结果
            for r in self.filtered_results:
                if (r['file_name'] == values[0] and
                    str(r['line_number']) == values[4] and
                    r['matched_text'] == values[3]):
                    selected.append(r)
                    break
        return selected

    def _open_file_location(self):
        """打开文件所在位置并选中文件"""
        selected = self._get_selected_results()
        if not selected:
            return
        file_path = selected[0]['file_path']
        if not os.path.exists(file_path):
            messagebox.showerror("错误", f"文件不存在:\n{file_path}")
            return
        import subprocess
        if sys.platform == 'win32':
            # Windows: explorer /select 打开文件夹并选中文件
            subprocess.run(['explorer', '/select,', file_path])
        elif sys.platform == 'darwin':
            # macOS: open -R 在Finder中定位文件
            subprocess.run(['open', '-R', file_path])
        else:
            # Linux: 打开所在目录
            subprocess.run(['xdg-open', os.path.dirname(file_path)])
        self.status_label.config(text=f"已定位: {os.path.basename(file_path)}")

    def _open_selected_file(self):
        """用系统默认程序打开选中的文件"""
        selected = self._get_selected_results()
        if not selected:
            return
        file_path = selected[0]['file_path']
        if not os.path.exists(file_path):
            messagebox.showerror("错误", f"文件不存在:\n{file_path}")
            return
        import subprocess
        if sys.platform == 'win32':
            os.startfile(file_path)
        elif sys.platform == 'darwin':
            subprocess.run(['open', file_path])
        else:
            subprocess.run(['xdg-open', file_path])
        self.status_label.config(text=f"已打开: {os.path.basename(file_path)}")

    def _on_double_click(self, event):
        """双击结果行时打开文件"""
        selection = self.tree.selection()
        if selection:
            self._open_selected_file()

    def _redact_selected(self):
        selected = self._get_selected_results()
        if not selected:
            messagebox.showinfo("提示", "请先选择要处理的条目")
            return

        count = len(selected)
        if not messagebox.askyesno("确认", f"将对 {count} 条匹配所在的文件进行脱敏处理（替换为***），\n原文件将备份为 .bak 文件。\n\n是否继续？"):
            return

        # 按文件分组去重
        files_to_process = {}
        for r in selected:
            fp = r['file_path']
            if fp not in files_to_process:
                # 找到该文件对应的匹配模式（取该文件的所有匹配）
                file_matches = [x for x in self.filtered_results if x['file_path'] == fp]
                # 用该文件匹配到的所有模式进行脱敏
                patterns = set(x['pattern_name'] for x in file_matches)
                patterns_info = [p for p in SENSITIVE_PATTERNS if p['name'] in patterns]
                files_to_process[fp] = patterns_info

        success_count = 0
        fail_count = 0
        results_detail = []

        for file_path, patterns in files_to_process.items():
            # 依次应用每个模式
            for pattern_info in patterns:
                success, msg = redact_file(file_path, pattern_info['pattern'], pattern_info['name'])
                if success:
                    success_count += 1
                    results_detail.append(f"✓ {os.path.basename(file_path)} - {pattern_info['name']}: {msg}")
                else:
                    if msg != "未找到匹配内容":  # 跳过"未找到"的提示
                        fail_count += 1
                        results_detail.append(f"✗ {os.path.basename(file_path)}: {msg}")

        summary = f"脱敏处理完成！\n成功: {success_count} 个文件，失败: {fail_count} 个"
        messagebox.showinfo("处理结果", summary)
        self.status_label.config(text=summary)
        self._refresh_after_action()

    def _delete_selected(self):
        selected = self._get_selected_results()
        if not selected:
            messagebox.showinfo("提示", "请先选择要处理的条目")
            return

        # 按文件去重
        files_to_delete = list(set(r['file_path'] for r in selected))
        count = len(files_to_delete)

        if not messagebox.askyesno("⚠️ 危险操作确认",
            f"将删除以下 {count} 个文件（移入废纸篓）：\n\n" +
            "\n".join(f"  • {os.path.basename(f)}" for f in files_to_delete) +
            "\n\n此操作不可撤销！是否继续？"):
            return

        # 再次确认
        if not messagebox.askyesno("二次确认", "⚠️ 删除后可通过废纸篓恢复，确定要删除吗？"):
            return

        success_count = 0
        fail_count = 0
        for file_path in files_to_delete:
            success, msg = safe_delete_file(file_path)
            if success:
                success_count += 1
            else:
                fail_count += 1
                messagebox.showerror("删除失败", f"{os.path.basename(file_path)}: {msg}")

        summary = f"删除完成！成功: {success_count} 个，失败: {fail_count} 个"
        messagebox.showinfo("处理结果", summary)
        self.status_label.config(text=summary)
        self._refresh_after_action()

    def _copy_match_content(self):
        selected = self._get_selected_results()
        if not selected:
            return
        self.root.clipboard_clear()
        text_to_copy = selected[0]['matched_text']
        self.root.clipboard_append(text_to_copy)
        self.status_label.config(text=f"已复制: {text_to_copy[:30]}...")

    def _refresh_after_action(self):
        """操作后刷新显示"""
        # 从结果中移除已不存在文件的匹配
        self.scan_results = [r for r in self.scan_results if os.path.exists(r['file_path'])]
        self._apply_filter()

        unique_files = len(set(r['file_path'] for r in self.scan_results))
        self.stats_label.config(
            text=f"已扫描: {self.engine.scanned_files} 个文件 | 发现: {unique_files} 个文件有风险 | 共计: {len(self.scan_results)} 条匹配"
        )

    def _export_report(self):
        """导出扫描报告为JSON/文本"""
        if not self.scan_results:
            messagebox.showinfo("提示", "没有扫描结果可以导出")
            return

        file_path = filedialog.asksaveasfilename(
            defaultextension=".txt",
            filetypes=[("文本文档", "*.txt"), ("JSON文件", "*.json")],
            title="导出报告"
        )
        if not file_path:
            return

        try:
            # 按文件分组整理报告
            report_data = {
                "scan_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "scan_directories": self.scan_dirs,
                "total_files_scanned": self.engine.scanned_files,
                "total_matches": len(self.scan_results),
                "risk_summary": {
                    "高": len([r for r in self.scan_results if r['risk'] == '高']),
                    "中": len([r for r in self.scan_results if r['risk'] == '中']),
                    "低": len([r for r in self.scan_results if r['risk'] == '低']),
                },
                "details": sorted(self.scan_results, key=lambda x: x['file_path'])
            }

            if file_path.endswith('.json'):
                with open(file_path, 'w', encoding='utf-8') as f:
                    json.dump(report_data, f, ensure_ascii=False, indent=2)
            else:
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.write("=" * 56 + "\n")
                    f.write("  ◆  数据安全检查报告  ◆\n")
                    f.write(" 广东省·开平农商银行·数据银行部 钻石专属\n")
                    f.write("=" * 56 + "\n\n")
                    f.write(f"  📅 扫描时间：{report_data['scan_time']}\n")
                    f.write(f"  📂 扫描目录：{', '.join(report_data['scan_directories'])}\n")
                    f.write(f"  📄 扫描文件：{report_data['total_files_scanned']} 个\n")
                    f.write(f"  🔍 匹配总数：{report_data['total_matches']} 条\n\n")
                    f.write(f"  ◆ 钻石专属 ◆  © 开平农商银行·数据银行部\n")
                    f.write("\n")

                    f.write("─" * 56 + "\n")
                    f.write("  风险统计\n")
                    f.write("─" * 56 + "\n")
                    f.write(f"    🔴 高风险：{report_data['risk_summary']['高']} 条\n")
                    f.write(f"    🟡 中风险：{report_data['risk_summary']['中']} 条\n")
                    f.write(f"    🟢 低风险：{report_data['risk_summary']['低']} 条\n\n")

                    f.write("─" * 56 + "\n")
                    f.write("  敏感信息详单\n")
                    f.write("─" * 56 + "\n\n")
                    for i, r in enumerate(report_data['details'], 1):
                        risk_symbol = "🔴" if r['risk'] == '高' else ("" if r['risk'] == '中' else "")
                        if r.get('cell_ref'):
                            location = f"{r['sheet_name']}!{r['cell_ref']} 单元格"
                        elif r.get('source_type') == 'page':
                            parts = str(r['line_number']).split('_l')
                            location = f"第{parts[0][1:]}页 第{parts[1]}行"
                        elif r.get('source_type') == 'slide':
                            location = f"第{r['line_number'][1:]}张幻灯片"
                        elif r.get('source_type') == 'database':
                            location = f"表{r.get('db_table','?')}.{r.get('db_column','?')} 第{r['line_number'][3:]}行"
                        elif isinstance(r['line_number'], int):
                            location = f"第 {r['line_number']} 行"
                        else:
                            location = str(r['line_number'])
                        f.write(f"  【第 {i} 条】{risk_symbol} [{r['risk']}] {r['pattern_name']}\n")
                        f.write(f"     文件：{r['file_path']} → {location}\n")
                        f.write(f"     内容：{r['matched_text']}\n")
                        f.write(f"     上下文：{r['context']}\n")
                        f.write("  " + "·" * 52 + "\n\n")

            messagebox.showinfo("导出成功", f"报告已保存到:\n{file_path}")
        except Exception as e:
            messagebox.showerror("导出失败", str(e))


# ============================================================
# 入口
# ============================================================
def main():
    root = tk.Tk()
    app = DataSafetyScannerApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
